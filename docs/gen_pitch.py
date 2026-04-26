"""
Crow Hub v2 — Investor / Partner Pitch Deck Generator
Generates a 10-slide professional presentation using python-pptx.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──────────────────────────────────────────────────
BG_DARK    = RGBColor(0x0D, 0x11, 0x17)   # slide background
BG_CARD    = RGBColor(0x16, 0x1B, 0x22)   # card / section
ACCENT     = RGBColor(0x58, 0xA6, 0xFF)   # primary blue
ACCENT2    = RGBColor(0x3B, 0x82, 0xF6)   # deeper blue
GREEN      = RGBColor(0x7E, 0xE7, 0x87)   # highlight green
ORANGE     = RGBColor(0xF0, 0x88, 0x3E)   # call-to-action
RED        = RGBColor(0xF8, 0x53, 0x49)   # problem / alert
PURPLE     = RGBColor(0xD2, 0xA8, 0xFF)   # memory / innovation
TEXT_WHITE = RGBColor(0xE6, 0xED, 0xF3)   # body text
TEXT_DIM   = RGBColor(0x8B, 0x94, 0x9E)   # dimmed / caption
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# ── Helpers ─────────────────────────────────────────────────────────

def dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_text(slide, text, x, y, w, h, size=16, color=TEXT_WHITE,
             bold=False, italic=False, align=PP_ALIGN.LEFT,
             font_name="Calibri", valign=MSO_ANCHOR.TOP, spacing=None):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = align
    if spacing:
        p.space_after = Pt(spacing)
    return txBox

def add_multiline(slide, lines, x, y, w, h, size=14, color=TEXT_WHITE,
                  bold=False, font_name="Calibri", line_spacing=1.3, bullet=False):
    """Add multiple lines of text. Each line is a string or (text, color, bold) tuple."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            txt, clr, bld = line
        else:
            txt, clr, bld = line, color, bold
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(size)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = font_name
        p.space_after = Pt(size * 0.5)
        p.line_spacing = Pt(size * line_spacing)
    return txBox

def add_rect(slide, x, y, w, h, color=BG_CARD, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, x, y, w, h, color=BG_CARD, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_accent_bar(slide, x, y, h, color=ACCENT, width=0.06):
    return add_rect(slide, x, y, width, h, color=color)

def card_with_accent(slide, x, y, w, h, accent_color=ACCENT):
    add_rect(slide, x, y, w, h, color=BG_CARD, border=RGBColor(0x30, 0x36, 0x3D))
    add_accent_bar(slide, x, y, h, color=accent_color)

def add_big_number(slide, number, label, x, y, num_color=GREEN):
    add_text(slide, number, x, y, 2.5, 0.7, size=42, color=num_color, bold=True, font_name="Calibri", align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y+0.65, 2.5, 0.5, size=12, color=TEXT_DIM, align=PP_ALIGN.CENTER)


# ====================================================================
# SLIDE 1 — TITLE
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
dark_bg(sl)

# Top accent line
add_rect(sl, 0, 0, 13.333, 0.04, color=ACCENT)

# Logo area — stylized "CH" monogram
add_rounded_rect(sl, 5.8, 1.2, 1.7, 1.7, color=ACCENT2)
add_text(sl, "CH", 5.8, 1.25, 1.7, 1.7, size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font_name="Consolas")

# Title
add_text(sl, "Crow Hub", 2.2, 3.3, 9, 0.9, size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font_name="Calibri")
add_text(sl, "Universal AI Agent Orchestration Platform", 2.2, 4.1, 9, 0.5, size=22, color=ACCENT, bold=False, align=PP_ALIGN.CENTER)

# Tagline
add_text(sl, "One hub.  Every agent.  Any model.  Shared memory.", 2.2, 5.0, 9, 0.5, size=16, color=TEXT_DIM, italic=True, align=PP_ALIGN.CENTER)

# Bottom bar
add_rect(sl, 0, 7.1, 13.333, 0.4, color=BG_CARD)
add_text(sl, "Built in Rust  |  Open Source  |  April 2026", 0.5, 7.12, 12, 0.35, size=11, color=TEXT_DIM, align=PP_ALIGN.CENTER)


# ====================================================================
# SLIDE 2 — THE PROBLEM
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)

add_accent_bar(sl, 0.7, 0.5, 0.5, color=RED)
add_text(sl, "The Problem", 0.9, 0.5, 5, 0.5, size=32, color=WHITE, bold=True)
add_text(sl, "AI development is fragmented", 0.9, 1.05, 8, 0.35, size=16, color=RED, italic=True)

# Left column — text
add_multiline(sl, [
    ("Developers juggle 5-10 agent interfaces across different", TEXT_WHITE, False),
    ("machines, operating systems, and hardware architectures.", TEXT_WHITE, False),
    ("", TEXT_WHITE, False),
    ("Each agent has its own model connection, its own memory", TEXT_WHITE, False),
    ("silo, its own communication protocol.", TEXT_WHITE, False),
    ("", TEXT_WHITE, False),
    ("No interop. No shared context. No unified orchestration.", ORANGE, True),
    ("", TEXT_WHITE, False),
    ("This chaos is NORMAL for serious AI developers", TEXT_DIM, True),
    ("and enterprises evaluating models across hardware.", TEXT_DIM, True),
], 0.9, 1.7, 5.8, 3.5, size=14)

# Right column — Real environment example
card_with_accent(sl, 7.2, 1.3, 5.5, 5.2, accent_color=ORANGE)
add_text(sl, "REAL TEST ENVIRONMENT", 7.5, 1.4, 5, 0.3, size=10, color=ORANGE, bold=True, font_name="Consolas")

env_lines = [
    ("Windows x86 PC", ACCENT, True),
    ("  LM Studio, Claude Desktop", TEXT_DIM, False),
    ("", TEXT_WHITE, False),
    ("WSL2 - Ubuntu Distro", ACCENT, True),
    ("  Claude Code, Gemini, Kimi, OpenClaw", TEXT_DIM, False),
    ("", TEXT_WHITE, False),
    ("WSL2 - NVIDIA Workbench Distro", ACCENT, True),
    ("  GPU compute, CUDA workloads", TEXT_DIM, False),
    ("", TEXT_WHITE, False),
    ("ARM64 DGX Spark  (SSH / Ethernet)", ACCENT, True),
    ("  vLLM, Ollama, Hermes", TEXT_DIM, False),
    ("  Claude Code, Gemini, Kimi, OpenClaw", TEXT_DIM, False),
]
add_multiline(sl, env_lines, 7.5, 1.85, 5, 4.5, size=12, font_name="Consolas", line_spacing=1.15)

# Bottom stat strip
add_rect(sl, 0.7, 5.7, 12, 1.1, color=BG_CARD, border=RGBColor(0x30, 0x36, 0x3D))
add_big_number(sl, "8+", "Agent interfaces", 0.9, 5.8, num_color=RED)
add_big_number(sl, "4", "Operating systems", 3.5, 5.8, num_color=RED)
add_big_number(sl, "2", "CPU architectures", 6.1, 5.8, num_color=RED)
add_big_number(sl, "0", "Interoperability", 8.7, 5.8, num_color=RED)
add_big_number(sl, "0", "Shared memory", 11.3, 5.8, num_color=RED)


# ====================================================================
# SLIDE 3 — THE VISION
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)

add_accent_bar(sl, 0.7, 0.5, 0.5, color=GREEN)
add_text(sl, "The Vision", 0.9, 0.5, 5, 0.5, size=32, color=WHITE, bold=True)

add_text(sl, "A high-performance Rust middleware that sits between\nALL your agents, ALL your models, and ALL your hardware.", 0.9, 1.2, 11, 0.8, size=18, color=TEXT_WHITE)

# Four pillars
pillars = [
    ("Hot-Swap\nAgents", "Plug in Claude, Gemini,\nCodex, Hermes, Kimi —\nany agent, anytime", ACCENT),
    ("Route Any\nModel", "vLLM, Ollama, LM Studio\nto any agent interface.\nAuto-discovery built in", GREEN),
    ("Shared\nMemory", "Graph-based knowledge\nthat all agents read\nand write together", PURPLE),
    ("Monitor\nEverything", "Tokens, cost, latency\nacross every agent\nand model backend", ORANGE),
]
for i, (title, desc, color) in enumerate(pillars):
    x = 0.7 + i * 3.1
    card_with_accent(sl, x, 2.4, 2.8, 3.6, accent_color=color)
    add_text(sl, title, x + 0.25, 2.6, 2.4, 0.9, size=20, color=color, bold=True)
    add_text(sl, desc, x + 0.25, 3.55, 2.3, 2.0, size=13, color=TEXT_DIM)

# Tagline
add_text(sl, "One interface to orchestrate your entire AI infrastructure.", 0.9, 6.3, 11.5, 0.5, size=15, color=TEXT_DIM, italic=True, align=PP_ALIGN.CENTER)


# ====================================================================
# SLIDE 4 — ARCHITECTURE OVERVIEW
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)

add_accent_bar(sl, 0.7, 0.5, 0.5, color=ACCENT)
add_text(sl, "Four-Layer Architecture", 0.9, 0.5, 8, 0.5, size=32, color=WHITE, bold=True)

layers = [
    ("Layer 1", "Model Backend", "Auto-discovers local servers (vLLM, Ollama, LM Studio, llama.cpp).\nMock API server exposes any backend as OpenAI-compatible.\nRoute any model to any agent interface.", ACCENT, 1.3),
    ("Layer 2", "Agent Interface", "Hot-swappable plugins: Claude, Gemini, Codex, OpenClaw, Hermes,\nKimi, OpenCode, Copaw + custom. TOML manifest registration.", GREEN, 2.65),
    ("Layer 3", "Communication", "Group channels (#code-review), Direct Messages between agents.\nConfigurable monitoring ACLs: FULL / SUMMARY / NOTIFY / NONE.", ORANGE, 4.0),
    ("Layer 4", "Graph Memory", "Typed knowledge graph (TASK / SKILL / EVENT). Embedding compression\n(~75% token reduction). Dual-path retrieval with PageRank.\nCross-workspace memory sharing and merging.", PURPLE, 5.35),
]

for label, title, desc, color, y_pos in layers:
    # Layer card
    add_rect(sl, 0.7, y_pos, 11.9, 1.15, color=BG_CARD, border=RGBColor(0x30, 0x36, 0x3D))
    add_accent_bar(sl, 0.7, y_pos, 1.15, color=color)

    # Layer number badge
    add_rounded_rect(sl, 1.0, y_pos + 0.15, 1.3, 0.4, color=color)
    add_text(sl, label, 1.0, y_pos + 0.15, 1.3, 0.4, size=12, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER, font_name="Consolas")

    # Title + description
    add_text(sl, title, 2.5, y_pos + 0.08, 5, 0.4, size=18, color=color, bold=True)
    add_text(sl, desc, 2.5, y_pos + 0.48, 9.8, 0.65, size=11, color=TEXT_DIM, font_name="Calibri")


# ====================================================================
# SLIDE 5 — CROSS-PLATFORM REALITY
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)

add_accent_bar(sl, 0.7, 0.5, 0.5, color=ORANGE)
add_text(sl, "Cross-Platform Reality", 0.9, 0.5, 8, 0.5, size=32, color=WHITE, bold=True)
add_text(sl, "Crow Hub orchestrates ALL of these from one interface", 0.9, 1.05, 10, 0.35, size=15, color=ORANGE, italic=True)

# Platform cards — 2x2 grid
platforms = [
    ("Windows x86 PC", ["LM Studio  (local models)", "Claude Desktop  (native app)"], ACCENT, 0.7, 1.7),
    ("WSL2 Ubuntu Distro", ["Claude Code", "Gemini, Kimi, OpenClaw"], GREEN, 6.85, 1.7),
    ("WSL2 NVIDIA Workbench", ["GPU compute workloads", "CUDA acceleration"], PURPLE, 0.7, 3.7),
    ("ARM64 DGX Spark", ["vLLM, Ollama  (model servers)", "Hermes, Claude Code, Gemini, Kimi, OpenClaw"], ORANGE, 6.85, 3.7),
]

for title, items, color, x, y in platforms:
    card_with_accent(sl, x, y, 5.8, 1.7, accent_color=color)
    add_text(sl, title, x + 0.25, y + 0.12, 5.3, 0.35, size=16, color=color, bold=True)
    lines = [(f"  {item}", TEXT_DIM, False) for item in items]
    add_multiline(sl, lines, x + 0.25, y + 0.55, 5.2, 1.1, size=12, font_name="Consolas", line_spacing=1.2)

# Connection diagram text
add_rect(sl, 3.5, 5.7, 6.3, 1.0, color=BG_CARD, border=ACCENT)
add_text(sl, "CROW HUB", 3.5, 5.7, 6.3, 0.45, size=22, color=ACCENT, bold=True, align=PP_ALIGN.CENTER, font_name="Consolas")
add_text(sl, "Unified orchestration  |  Shared memory  |  Model routing", 3.5, 6.15, 6.3, 0.4, size=12, color=TEXT_DIM, align=PP_ALIGN.CENTER)


# ====================================================================
# SLIDE 6 — WHY RUST
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)

add_accent_bar(sl, 0.7, 0.5, 0.5, color=ORANGE)
add_text(sl, "Why Rust", 0.9, 0.5, 5, 0.5, size=32, color=WHITE, bold=True)
add_text(sl, "The right tool for high-performance multi-agent orchestration", 0.9, 1.05, 10, 0.35, size=15, color=ORANGE, italic=True)

reasons = [
    ("Safety", "Memory-safe concurrency without garbage collection.\nFearless multi-threaded agent orchestration.\nNo data races, no null pointers, no segfaults.", ACCENT),
    ("Performance", "Zero-cost abstractions. Native speed.\nAsync I/O via Tokio runtime.\nDashMap lock-free concurrent maps.", GREEN),
    ("Cross-Platform", "Compiles natively to Windows, Linux, ARM64.\nSingle binary deployment — no runtime deps.\nRuns on x86 PCs and DGX Spark alike.", ORANGE),
    ("Ecosystem", "Tokio  — production async runtime\nRatatui  — rich terminal UI\nTauri  — cross-platform desktop GUI\nSerde  — zero-copy serialization", PURPLE),
]

for i, (title, desc, color) in enumerate(reasons):
    x = 0.7 + (i % 2) * 6.15
    y = 1.7 + (i // 2) * 2.6
    card_with_accent(sl, x, y, 5.9, 2.3, accent_color=color)
    add_text(sl, title, x + 0.25, y + 0.15, 5.4, 0.4, size=22, color=color, bold=True)
    add_text(sl, desc, x + 0.25, y + 0.65, 5.4, 1.5, size=13, color=TEXT_DIM)


# ====================================================================
# SLIDE 7 — GRAPH MEMORY INNOVATION
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)

add_accent_bar(sl, 0.7, 0.5, 0.5, color=PURPLE)
add_text(sl, "Graph Memory Innovation", 0.9, 0.5, 8, 0.5, size=32, color=WHITE, bold=True)
add_text(sl, "Inspired by graph-memory  —  adapted for multi-agent orchestration", 0.9, 1.05, 11, 0.35, size=15, color=PURPLE, italic=True)

# Left column — knowledge graph
card_with_accent(sl, 0.7, 1.6, 5.8, 2.4, accent_color=PURPLE)
add_text(sl, "Knowledge Graph", 0.95, 1.7, 5.3, 0.35, size=18, color=PURPLE, bold=True)
graph_lines = [
    ("TASK nodes  — completed actions", GREEN, False),
    ("SKILL nodes  — procedural know-how", ACCENT, False),
    ("EVENT nodes  — failures & incidents", RED, False),
    ("", TEXT_WHITE, False),
    ("Edges: USED_SKILL, SOLVED_BY,", TEXT_DIM, False),
    ("REQUIRES, PATCHES, CONFLICTS_WITH", TEXT_DIM, False),
]
add_multiline(sl, graph_lines, 0.95, 2.15, 5.3, 1.8, size=12, font_name="Consolas", line_spacing=1.25)

# Right column — compression stats
card_with_accent(sl, 6.85, 1.6, 5.8, 2.4, accent_color=GREEN)
add_text(sl, "Embedding Compression", 7.1, 1.7, 5.3, 0.35, size=18, color=GREEN, bold=True)
comp_lines = [
    ("Raw messages -> structured triples", TEXT_WHITE, False),
    ("~75% token reduction", GREEN, True),
    ("Cosine dedup (>= 0.90 -> merge)", TEXT_DIM, False),
    ("Community summarization via LLM", TEXT_DIM, False),
    ("", TEXT_WHITE, False),
    ("174 msgs = 95K tokens -> ~24K", ORANGE, True),
]
add_multiline(sl, comp_lines, 7.1, 2.15, 5.3, 1.8, size=12, font_name="Consolas", line_spacing=1.25)

# Bottom row — dual retrieval + portability
card_with_accent(sl, 0.7, 4.25, 5.8, 2.5, accent_color=ACCENT)
add_text(sl, "Dual-Path Retrieval", 0.95, 4.35, 5.3, 0.35, size=18, color=ACCENT, bold=True)
retrieval_lines = [
    ("Precise:  vector/FTS5 -> seed -> PageRank", TEXT_WHITE, False),
    ("General:  community embeddings -> PPR", TEXT_WHITE, False),
    ("", TEXT_WHITE, False),
    ("Personalized PageRank ranks by relevance", GREEN, False),
    ("to YOUR current task, not just recency", GREEN, False),
    ("", TEXT_WHITE, False),
    ("Community detection clusters related skills", TEXT_DIM, False),
]
add_multiline(sl, retrieval_lines, 0.95, 4.8, 5.3, 1.9, size=12, font_name="Consolas", line_spacing=1.15)

card_with_accent(sl, 6.85, 4.25, 5.8, 2.5, accent_color=ORANGE)
add_text(sl, "Portable & Shareable", 7.1, 4.35, 5.3, 0.35, size=18, color=ORANGE, bold=True)
portable_lines = [
    ("Standard format across ALL agents", TEXT_WHITE, False),
    ("Agent-specific memory plugins translate", TEXT_WHITE, False),
    ("each agent's context to shared graph", TEXT_WHITE, False),
    ("", TEXT_WHITE, False),
    ("Export .crow-memory bundles", GREEN, False),
    ("Import & merge across workspaces", GREEN, False),
    ("Cosine dedup prevents duplicates", TEXT_DIM, False),
]
add_multiline(sl, portable_lines, 7.1, 4.8, 5.3, 1.9, size=12, font_name="Consolas", line_spacing=1.15)


# ====================================================================
# SLIDE 8 — MARKET OPPORTUNITY
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)

add_accent_bar(sl, 0.7, 0.5, 0.5, color=GREEN)
add_text(sl, "Market Opportunity", 0.9, 0.5, 8, 0.5, size=32, color=WHITE, bold=True)

# Big stats row
add_rect(sl, 0.7, 1.3, 12, 1.5, color=BG_CARD, border=RGBColor(0x30, 0x36, 0x3D))
add_big_number(sl, "40%+", "CAGR agent market", 1.2, 1.4, num_color=GREEN)
add_big_number(sl, "$28B", "AI agent TAM by 2028", 4.0, 1.4, num_color=GREEN)
add_big_number(sl, "0", "Universal orchestration\nstandards exist", 6.8, 1.4, num_color=RED)
add_big_number(sl, "100%", "Developers building\ncustom glue code", 9.6, 1.4, num_color=ORANGE)

# Opportunity cards
opps = [
    ("Enterprise Evaluation Gap", "Every enterprise testing LLMs needs to compare models across\nhardware (x86, ARM, GPU). No tool does this seamlessly today.\nCrow Hub is the first unified evaluation + orchestration layer.", ACCENT),
    ("Developer Tooling Vacuum", "AI developers spend 30%+ of their time on integration glue.\nEvery team reinvents agent communication and memory.\nCrow Hub replaces thousands of lines of bespoke glue code.", GREEN),
    ("Open-Source Adoption Flywheel", "Community-driven plugin ecosystem. Each new agent plugin\nmakes the platform more valuable. Network effects compound\nas users share memory graphs and agent configurations.", ORANGE),
]

for i, (title, desc, color) in enumerate(opps):
    y = 3.1 + i * 1.4
    card_with_accent(sl, 0.7, y, 12, 1.2, accent_color=color)
    add_text(sl, title, 0.95, y + 0.1, 4, 0.35, size=16, color=color, bold=True)
    add_text(sl, desc, 5.0, y + 0.1, 7.5, 1.0, size=12, color=TEXT_DIM)


# ====================================================================
# SLIDE 9 — DEVELOPMENT ROADMAP
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)

add_accent_bar(sl, 0.7, 0.5, 0.5, color=ACCENT)
add_text(sl, "Development Roadmap", 0.9, 0.5, 8, 0.5, size=32, color=WHITE, bold=True)
add_text(sl, "28 weeks  |  8 phases  |  Incremental delivery", 0.9, 1.05, 8, 0.35, size=15, color=ACCENT, italic=True)

phases = [
    ("Phase 0-1", "Foundation + Core", "5 weeks", "Protocol, Hub runtime,\nComms layer, ACLs", ACCENT, 0.7, 1.7),
    ("Phase 2", "Agent Plugin System", "4 weeks", "Plugin loader, TOML manifests\n3 initial agents", GREEN, 4.0, 1.7),
    ("Phase 3", "Graph Memory Engine", "4 weeks", "Knowledge graph, embeddings,\nPPR, dual-path retrieval", PURPLE, 7.3, 1.7),
    ("Phase 3b", "Model Router", "3 weeks", "Model registry, auto-discovery\nMock API server", ORANGE, 10.6, 1.7),
    ("Phase 4", "Monitor + Cost", "2 weeks", "Per-agent metrics,\ncost tracking, exporters", ACCENT, 0.7, 4.0),
    ("Phase 5", "TUI Dashboards", "3 weeks", "Channel views, memory browser\nmodel routing dashboard", GREEN, 4.0, 4.0),
    ("Phase 6", "Memory Sharing", "2 weeks", "Memory plugins, export/import\ncross-workspace merge", PURPLE, 7.3, 4.0),
    ("Phase 7-8", "GUI + Release", "5 weeks", "Tauri desktop GUI,\ntesting, more agents", ORANGE, 10.6, 4.0),
]

for label, title, duration, desc, color, x, y in phases:
    card_with_accent(sl, x, y, 3.0, 1.95, accent_color=color)
    add_rounded_rect(sl, x + 0.2, y + 0.12, 1.2, 0.3, color=color)
    add_text(sl, label, x + 0.2, y + 0.12, 1.2, 0.3, size=9, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER, font_name="Consolas")
    add_text(sl, duration, x + 1.55, y + 0.12, 1.2, 0.3, size=10, color=color, bold=True, font_name="Consolas")
    add_text(sl, title, x + 0.2, y + 0.55, 2.6, 0.35, size=14, color=WHITE, bold=True)
    add_text(sl, desc, x + 0.2, y + 0.95, 2.6, 0.9, size=11, color=TEXT_DIM)

# Timeline bar at bottom
add_rect(sl, 0.7, 6.3, 12, 0.06, color=ACCENT)
months = ["Apr '26", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov '26"]
for i, m in enumerate(months):
    x_pos = 0.7 + i * (12.0/7)
    add_text(sl, m, x_pos - 0.4, 6.45, 1.2, 0.3, size=10, color=TEXT_DIM, align=PP_ALIGN.CENTER, font_name="Consolas")
    # Tick mark
    add_rect(sl, x_pos + 0.15, 6.22, 0.02, 0.15, color=ACCENT)


# ====================================================================
# SLIDE 10 — THE ASK
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)

# Top accent
add_rect(sl, 0, 0, 13.333, 0.04, color=ACCENT)

# "CH" monogram again for brand consistency
add_rounded_rect(sl, 5.8, 0.7, 1.7, 1.7, color=ACCENT2)
add_text(sl, "CH", 5.8, 0.75, 1.7, 1.7, size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font_name="Consolas")

# Main CTA
add_text(sl, "We're building the middleware layer that\nevery AI developer needs but nobody has built yet.",
         1.5, 2.7, 10.3, 1.2, size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(sl, "Join us.", 1.5, 3.95, 10.3, 0.6, size=28, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# Role cards
roles = [
    ("Systems Engineers", "Rust / C++ experience\nAsync systems, networking", ACCENT),
    ("ML Infrastructure", "Model serving, embeddings\nvLLM, Ollama, CUDA", GREEN),
    ("Open Source Contributors", "Agent plugins, memory plugins\nTesting, documentation", ORANGE),
]

for i, (title, desc, color) in enumerate(roles):
    x = 1.2 + i * 3.8
    card_with_accent(sl, x, 4.8, 3.5, 1.6, accent_color=color)
    add_text(sl, title, x + 0.25, 4.9, 3.0, 0.35, size=16, color=color, bold=True)
    add_text(sl, desc, x + 0.25, 5.35, 3.0, 0.9, size=12, color=TEXT_DIM)

# Bottom bar
add_rect(sl, 0, 6.9, 13.333, 0.6, color=BG_CARD)
add_text(sl, "github.com/crow-hub  |  Apache 2.0  |  contact@crowhub.dev", 0.5, 6.95, 12.3, 0.45, size=13, color=TEXT_DIM, align=PP_ALIGN.CENTER)


# ── SAVE ───────────────────────────────────────────────────────────
out_path = r"D:\CodeBuddy\20260410160026\crow-hub\docs\crow-hub-pitch.pptx"
prs.save(out_path)
print(f"Saved to {out_path}")
print(f"Slides: {len(prs.slides)}")
